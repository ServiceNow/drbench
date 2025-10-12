"""
Enhanced content processing system for DrBench Agent
Handles file downloads, text extraction, and vector store integration
"""

import mimetypes
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup
from drbench.agents.utils import prompt_llm

from drbench.agents.drbench_agent.session_cache import SessionCache
from drbench.agents.drbench_agent.vector_store import VectorStore


class ContentProcessor:
    """Comprehensive content processing for files and URLs"""

    def __init__(
        self,
        workspace_dir: str,
        model: str,
        vector_store: VectorStore = None,
        session_cache: Optional[SessionCache] = None,
    ):
        self.workspace_dir = Path(workspace_dir)
        self.workspace_dir.mkdir(exist_ok=True)
        self.vector_store = vector_store
        self.model = model
        self.session_cache = session_cache

        # Create subdirectories
        self.downloads_dir = self.workspace_dir / "downloads"
        self.extracted_dir = self.workspace_dir / "extracted_content"
        self.downloads_dir.mkdir(exist_ok=True)
        self.extracted_dir.mkdir(exist_ok=True)

        # Session for web requests
        self.session = requests.Session()
        self.session.headers.update(
            {"User-Agent": "Mozilla/5.0 (compatible; DrBench-Agent/1.0; Research-Bot)", "Accept": "*/*"}
        )

    def process_url(self, url: str, query_context: str = "") -> Dict[str, Any]:
        """
        Download and process content from a URL

        Args:
            url: URL to process
            query_context: Context about what we're researching

        Returns:
            Processing result with extracted content and file paths
        """
        try:
            # Download content
            response = self.session.get(url, timeout=30)
            response.raise_for_status()

            # Determine content type and file extension
            content_type = response.headers.get("content-type", "").lower()
            parsed_url = urlparse(url)

            # Generate filename
            if parsed_url.path:
                filename = Path(parsed_url.path).name
                if not filename or "." not in filename:
                    filename = f"content_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            else:
                filename = f"content_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

            # Add appropriate extension based on content type
            if "text/html" in content_type and not filename.endswith(".html"):
                filename += ".html"
            elif "application/pdf" in content_type and not filename.endswith(".pdf"):
                filename += ".pdf"
            elif "application/json" in content_type and not filename.endswith(".json"):
                filename += ".json"
            elif "text/plain" in content_type and not filename.endswith(".txt"):
                filename += ".txt"

            # Save to downloads directory
            file_path = self.downloads_dir / filename
            with open(file_path, "wb") as f:
                f.write(response.content)

            # Extract text content
            extracted_content = self.extract_text_from_file(file_path, content_type)

            # Save extracted text
            text_filename = f"{Path(filename).stem}_extracted.txt"
            text_path = self.extracted_dir / text_filename
            with open(text_path, "w", encoding="utf-8") as f:
                f.write(extracted_content)

            # Check cache first
            doc_id = None
            cached = False

            if self.session_cache:
                # Check if URL has been processed
                cached_doc_id = self.session_cache.check_source("url", url)
                if cached_doc_id:
                    doc_id = cached_doc_id
                    cached = True
                    # Update context in both session cache and vector store
                    if query_context:
                        self.session_cache.add_document(
                            doc_id,
                            extracted_content,
                            source_type="url",
                            source_identifier=url,
                            query_context=query_context,
                        )
                        # Also update vector store metadata
                        if self.vector_store:
                            metadata = {
                                "source": "url",
                                "url": url,
                                "source_identifier": url,
                                "query_context": query_context,
                                "timestamp": datetime.now().isoformat(),
                            }
                            self.vector_store._merge_metadata(doc_id, metadata)

            # Store in vector store if available and not cached
            if not cached and self.vector_store and extracted_content.strip():
                metadata = {
                    "source": "url",
                    "url": url,
                    "source_identifier": url,  # Add for deduplication
                    "content_type": content_type,
                    "filename": filename,
                    "query_context": query_context,
                    "timestamp": datetime.now().isoformat(),
                    "file_path": str(file_path),
                    "extracted_path": str(text_path),
                }
                doc_id = self.vector_store.store_document(
                    content=extracted_content,
                    metadata=metadata,
                )

                # Add to cache
                if self.session_cache and doc_id:
                    self.session_cache.add_document(
                        doc_id,
                        extracted_content,
                        source_type="url",
                        source_identifier=url,
                        file_path=str(file_path),
                        query_context=query_context,
                    )

            return {
                "success": True,
                "url": url,
                "file_path": str(file_path),
                "extracted_path": str(text_path),
                "content_type": content_type,
                "content_length": len(extracted_content),
                "doc_id": doc_id,
                "stored_in_vector": doc_id is not None,
            }

        except Exception as e:
            return {"success": False, "url": url, "error": str(e)}

    def process_file(
        self, file_path: str, query_context: str = "", additional_metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Process an existing file (extract text and store in vector store)

        Args:
            file_path: Path to file to process
            query_context: Context about what we're researching
            additional_metadata: Additional metadata to store with the document

        Returns:
            Processing result
        """
        try:
            file_path = Path(file_path)

            if not file_path.exists():
                return {"success": False, "file_path": str(file_path), "error": "File does not exist"}

            # Determine content type
            content_type, _ = mimetypes.guess_type(str(file_path))
            if not content_type:
                content_type = "application/octet-stream"

            # Extract text content
            extracted_content = self.extract_text_from_file(file_path, content_type)

            # Save extracted text
            text_filename = f"{file_path.stem}_extracted.txt"
            text_path = self.extracted_dir / text_filename
            with open(text_path, "w", encoding="utf-8") as f:
                f.write(extracted_content)

            # Check cache first
            doc_id = None
            cached = False

            # Determine source identifier
            source_identifier = str(file_path)
            if additional_metadata:
                # Use service-specific identifier if available
                source_identifier = (
                    additional_metadata.get("source_identifier")
                    or additional_metadata.get("original_path")
                    or str(file_path)
                )

            if self.session_cache:
                # Check if file has been processed
                cached_result = self.session_cache.check_file(str(file_path))
                if cached_result:
                    _, doc_id = cached_result
                    cached = True
                else:
                    # Check by source identifier
                    source_type = additional_metadata.get("source", "file") if additional_metadata else "file"
                    cached_doc_id = self.session_cache.check_source(source_type, source_identifier)
                    if cached_doc_id:
                        doc_id = cached_doc_id
                        cached = True

                # Update context if cached
                if cached and query_context:
                    self.session_cache.add_document(
                        doc_id,
                        extracted_content,
                        source_type=additional_metadata.get("source", "file") if additional_metadata else "file",
                        source_identifier=source_identifier,
                        file_path=str(file_path),
                        query_context=query_context,
                    )
                    # Also update vector store metadata
                    if self.vector_store:
                        metadata = {
                            "source": additional_metadata.get("source", "file") if additional_metadata else "file",
                            "source_identifier": source_identifier,
                            "query_context": query_context,
                            "timestamp": datetime.now().isoformat(),
                        }
                        if additional_metadata:
                            metadata.update(additional_metadata)
                        self.vector_store._merge_metadata(doc_id, metadata)

            # Store in vector store if available and not cached
            if not cached and self.vector_store and extracted_content.strip():
                # Base metadata
                metadata = {
                    "source": "file",
                    "original_path": str(file_path),
                    "source_identifier": source_identifier,  # Add for deduplication
                    "content_type": content_type,
                    "filename": file_path.name,
                    "query_context": query_context,
                    "timestamp": datetime.now().isoformat(),
                    "extracted_path": str(text_path),
                }

                # Add additional metadata if provided
                if additional_metadata:
                    metadata.update(additional_metadata)

                doc_id = self.vector_store.store_document(content=extracted_content, metadata=metadata)

                # Add to cache
                if self.session_cache and doc_id:
                    self.session_cache.add_document(
                        doc_id,
                        extracted_content,
                        source_type=metadata.get("source"),
                        source_identifier=source_identifier,
                        file_path=str(file_path),
                        query_context=query_context,
                    )

            return {
                "success": True,
                "file_path": str(file_path),
                "extracted_path": str(text_path),
                "content_type": content_type,
                "content_length": len(extracted_content),
                "doc_id": doc_id,
                "stored_in_vector": doc_id is not None,
            }

        except Exception as e:
            return {"success": False, "file_path": str(file_path), "error": str(e)}

    def extract_text_from_file(self, file_path: Path, content_type: str = None) -> str:
        """
        Extract text from various file formats

        Args:
            file_path: Path to file
            content_type: MIME type of file

        Returns:
            Extracted text content
        """
        if not content_type:
            content_type, _ = mimetypes.guess_type(str(file_path))

        try:
            if not content_type:
                content_type = "text/plain"

            content_type = content_type.lower()

            # Handle different file types
            if "text/html" in content_type:
                return self._extract_from_html(file_path)
            elif "application/pdf" in content_type:
                return self._extract_from_pdf(file_path)
            elif "application/json" in content_type:
                return self._extract_from_json(file_path)
            elif content_type.startswith("text/"):
                return self._extract_from_text(file_path)
            elif "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in content_type:
                return self._extract_from_docx(file_path)
            elif "application/msword" in content_type:
                return self._extract_from_doc(file_path)
            elif "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" in content_type:
                return self._extract_from_xlsx(file_path)
            elif "application/vnd.ms-excel" in content_type:
                return self._extract_from_xls(file_path)
            elif "application/vnd.openxmlformats-officedocument.presentationml.presentation" in content_type:
                return self._extract_from_pptx(file_path)
            elif "application/vnd.ms-powerpoint" in content_type:
                return self._extract_from_ppt(file_path)
            elif str(file_path).lower().endswith((".xlsx", ".xls")):
                # Handle Excel files by extension if MIME type detection fails
                if str(file_path).lower().endswith(".xlsx"):
                    return self._extract_from_xlsx(file_path)
                else:
                    return self._extract_from_xls(file_path)
            elif str(file_path).lower().endswith((".pptx", ".ppt")):
                # Handle PowerPoint files by extension if MIME type detection fails
                if str(file_path).lower().endswith(".pptx"):
                    return self._extract_from_pptx(file_path)
                else:
                    return self._extract_from_ppt(file_path)
            else:
                # Try as text file as fallback
                return self._extract_from_text(file_path)

        except Exception as e:
            return f"Error extracting text from {file_path}: {str(e)}"

    def _extract_from_html(self, file_path: Path) -> str:
        """Extract text from HTML file"""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()

        soup = BeautifulSoup(content, "html.parser")

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Extract text
        text = soup.get_text()

        # Clean up text
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        clean_text = " ".join(chunk for chunk in chunks if chunk)

        return clean_text

    def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file"""
        try:
            import PyPDF2

            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"

            return text.strip()

        except ImportError:
            try:
                import pdfplumber

                with pdfplumber.open(file_path) as pdf:
                    text = ""
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"

                return text.strip()

            except ImportError:
                return f"PDF processing libraries not available. Install PyPDF2 or pdfplumber to process PDF files. File: {file_path}"

    def _extract_from_json(self, file_path: Path) -> str:
        """Extract text from JSON file"""
        import json

        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        # Convert JSON to readable text
        json_str = json.dumps(data, indent=2, ensure_ascii=False)

        # Also try to extract any text values from the JSON structure
        text_content = []

        def extract_text_values(obj, path=""):
            if isinstance(obj, dict):
                for key, value in obj.items():
                    new_path = f"{path}.{key}" if path else key
                    if isinstance(value, str) and len(value) > 10:  # Meaningful text
                        text_content.append(f"{new_path}: {value}")
                    else:
                        extract_text_values(value, new_path)
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    extract_text_values(item, f"{path}[{i}]")
            elif isinstance(obj, str) and len(obj) > 10:
                text_content.append(f"{path}: {obj}")

        extract_text_values(data)

        # Combine JSON structure and extracted text
        if text_content:
            return f"JSON Structure:\n{json_str}\n\nExtracted Text Content:\n" + "\n".join(text_content)
        else:
            return json_str

    def _extract_from_text(self, file_path: Path) -> str:
        """Extract text from plain text file"""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    def _extract_from_docx(self, file_path: Path) -> str:
        """Extract text from Word DOCX file"""
        try:
            import docx

            doc = docx.Document(file_path)
            text = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)

            return "\n".join(text)

        except ImportError:
            return f"python-docx library not available. Install python-docx to process DOCX files. File: {file_path}"

    def _extract_from_doc(self, file_path: Path) -> str:
        """Extract text from Word DOC file"""
        try:
            import textract

            text = textract.process(str(file_path)).decode("utf-8")
            return text

        except ImportError:
            return f"textract library not available. Install textract to process DOC files. File: {file_path}"

    def _extract_from_xlsx(self, file_path: Path) -> str:
        """Extract text from Excel XLSX file"""
        try:
            import openpyxl

            workbook = openpyxl.load_workbook(file_path)
            text_content = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"Sheet: {sheet_name}")

                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text_content.append(" | ".join(row_text))

                text_content.append("")  # Empty line between sheets

            return "\n".join(text_content)

        except ImportError:
            return f"openpyxl library not available. Install openpyxl to process XLSX files. File: {file_path}"

    def _extract_from_xls(self, file_path: Path) -> str:
        """Extract text from Excel XLS file"""
        try:
            import xlrd

            workbook = xlrd.open_workbook(file_path)
            text_content = []

            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                text_content.append(f"Sheet: {sheet.name}")

                for row_idx in range(sheet.nrows):
                    row_text = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        if cell.value:
                            row_text.append(str(cell.value))
                    if row_text:
                        text_content.append(" | ".join(row_text))

                text_content.append("")  # Empty line between sheets

            return "\n".join(text_content)

        except ImportError:
            return f"xlrd library not available. Install xlrd to process XLS files. File: {file_path}"

    def _extract_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint PPTX file"""
        try:
            from pptx import Presentation

            presentation = Presentation(file_path)
            text_content = []

            for slide_idx, slide in enumerate(presentation.slides, 1):
                text_content.append(f"Slide {slide_idx}:")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content.append(shape.text)

                text_content.append("")  # Empty line between slides

            return "\n".join(text_content)

        except ImportError:
            return f"python-pptx library not available. Install python-pptx to process PPTX files. File: {file_path}"

    def _extract_from_ppt(self, file_path: Path) -> str:
        """Extract text from PowerPoint PPT file"""
        try:
            import textract

            text = textract.process(str(file_path)).decode("utf-8")
            return text

        except ImportError:
            return f"textract library not available. Install textract to process PPT files. File: {file_path}"

    def download_from_enterprise_service(
        self, service_result: Dict[str, Any], query_context: str = ""
    ) -> List[Dict[str, Any]]:
        """
        Download files identified from enterprise services

        Args:
            service_result: Result from enterprise API tool
            query_context: Context about the research

        Returns:
            List of processing results for downloaded files
        """
        results = []

        # Extract file URLs or paths from service result
        file_references = self._extract_file_references(service_result)

        for file_ref in file_references:
            if file_ref.get("type") == "url":
                result = self.process_url(file_ref["url"], query_context)
                results.append(result)
            elif file_ref.get("type") == "path":
                # For local file paths, just process them
                result = self.process_file(file_ref["path"], query_context)
                results.append(result)

        return results

    def _extract_file_references(self, service_result: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Extract file references from enterprise service results"""
        file_refs = []

        # Look for common patterns in service results
        if isinstance(service_result, dict):
            # Check for direct file URLs
            if "url" in service_result and isinstance(service_result["url"], str):
                file_refs.append({"type": "url", "url": service_result["url"]})

            # Check for file paths
            if "file_path" in service_result:
                file_refs.append({"type": "path", "path": service_result["file_path"]})

            # Check for arrays of files
            for key in ["files", "documents", "attachments", "results"]:
                if key in service_result and isinstance(service_result[key], list):
                    for item in service_result[key]:
                        if isinstance(item, dict):
                            if "url" in item:
                                file_refs.append({"type": "url", "url": item["url"]})
                            elif "path" in item:
                                file_refs.append({"type": "path", "path": item["path"]})
                            elif "file_path" in item:
                                file_refs.append({"type": "path", "path": item["file_path"]})

        return file_refs

    def chunk_content(self, content: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Split large content into chunks for better vector storage

        Args:
            content: Text content to chunk
            chunk_size: Maximum size of each chunk
            overlap: Number of characters to overlap between chunks

        Returns:
            List of content chunks
        """
        if len(content) <= chunk_size:
            return [content]

        chunks = []
        start = 0

        while start < len(content):
            end = start + chunk_size

            # Try to break at sentence or paragraph boundaries
            if end < len(content):
                # Look for sentence endings
                for i in range(end, max(start + chunk_size // 2, end - 200), -1):
                    if content[i] in ".!?":
                        end = i + 1
                        break
                # If no sentence break found, look for paragraph breaks
                else:
                    for i in range(end, max(start + chunk_size // 2, end - 200), -1):
                        if content[i : i + 2] == "\n\n":
                            end = i + 2
                            break

            chunk = content[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - overlap

            if start >= len(content):
                break

        return chunks

    def store_chunked_content(self, content: str, metadata: Dict[str, Any]) -> List[str]:
        """
        Store large content as chunks in vector store

        Args:
            content: Content to store
            metadata: Metadata for the content

        Returns:
            List of document IDs for stored chunks
        """
        if not self.vector_store:
            return []

        chunks = self.chunk_content(content)
        doc_ids = []

        for i, chunk in enumerate(chunks):
            chunk_metadata = metadata.copy()
            chunk_metadata.update({"chunk_index": i, "total_chunks": len(chunks), "is_chunk": True})

            doc_id = self.vector_store.store_document(content=chunk, metadata=chunk_metadata)
            doc_ids.append(doc_id)

            # Add to cache if first chunk
            if self.session_cache and i == 0:
                self.session_cache.add_document(
                    doc_id,
                    content,  # Store full content reference
                    source_type=metadata.get("source"),
                    source_identifier=metadata.get("source_identifier"),
                    file_path=metadata.get("file_path"),
                    query_context=metadata.get("query_context"),
                )

        return doc_ids

    def extract_sources_from_content(self, content: str, content_type: str = "document") -> Dict[str, Any]:
        """
        Extract citations and source references from content using LLM

        Args:
            content: Text content to analyze
            content_type: Type of content (document, message, etc.)

        Returns:
            Dict containing extracted sources and metadata
        """
        # Limit content length for LLM processing
        max_length = 8000
        if len(content) > max_length:
            # Take beginning and end to capture both intro citations and bibliography
            content_sample = content[: max_length // 2] + "\n...\n" + content[-max_length // 2 :]
        else:
            content_sample = content

        extraction_prompt = f"""
        Analyze the following {content_type} and extract all references to external sources, citations, and data sources.
        
        Content:
        {content_sample}
        
        Extract the following information:
        1. Academic citations (papers, journals, books)
        2. Reports and whitepapers
        3. Websites and online resources
        4. Databases and datasets
        5. News articles and media sources
        6. Internal documents or reports mentioned
        7. Statistics with their sources
        8. Quotes with attribution
        
        For each source found, provide:
        - source_type: (paper, report, website, database, news, internal_doc, statistic, quote)
        - title: The title or description
        - authors: Authors if mentioned (for papers/reports)
        - year: Publication year if mentioned
        - url: URL if provided
        - context: How it's referenced in the document
        
        Return as JSON with structure:
        {{
            "sources": [list of sources],
            "statistics": [list of statistics with their sources],
            "key_claims": [important claims that reference sources]
        }}
        
        If no sources are found, return an empty structure.
        """

        try:
            response = prompt_llm(model=self.model, prompt=extraction_prompt)

            if isinstance(response, dict):
                return response
            else:
                # Try to parse as JSON if string
                import json

                return json.loads(response)

        except Exception as e:
            return {"sources": [], "statistics": [], "key_claims": [], "extraction_error": str(e)}

    def get_stats(self) -> Dict[str, Any]:
        """Get statistics about processed content"""
        downloads_count = len(list(self.downloads_dir.glob("*"))) if self.downloads_dir.exists() else 0
        extracted_count = len(list(self.extracted_dir.glob("*"))) if self.extracted_dir.exists() else 0

        return {
            "downloads_directory": str(self.downloads_dir),
            "extracted_directory": str(self.extracted_dir),
            "downloaded_files": downloads_count,
            "extracted_files": extracted_count,
            "vector_store_connected": self.vector_store is not None,
        }
